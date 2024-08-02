
"""
generate_ppt_from_txt.py

Generate a PowerPoint presentation from a text file.

Copyright (c) 2023 Astra
Maintainer: Astra <astralee95@gmail.com>

Licensed under the MIT License. See the LICENSE file for more details.
"""

import os
import ctypes
import sys

# Set the console encoding to UTF-8
if os.name == 'nt':
    import msvcrt
    msvcrt.setmode(sys.stdout.fileno(), os.O_BINARY)
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

import argparse
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from lxml import etree
from PIL import Image

def rgb_to_ansi(r, g, b):
    """
    Convert RGB values to the closest ANSI color code.
    """
    # Calculate the closest 8-bit color code
    if r == g == b:
        if r < 8:
            return 16
        if r > 248:
            return 231
        return round(((r - 8) / 247) * 24) + 232

    return 16 + (36 * round(r / 255 * 5)) + (6 * round(g / 255 * 5)) + round(b / 255 * 5)

# Define the RGB color
rgb_color = (49, 51, 158)

# Convert RGB to ANSI escape code for background
ansi_bg_code = rgb_to_ansi(*rgb_color)
PPT_BG_BLUE = ''

# Define ANSI escape code for white text
WHITE_TEXT = '\033[38;5;15m'

# Define ANSI escape codes for bold and color
BOLD = '\033[1m'
ITALIC = '\033[3m'
RESET = '\033[0m'
RED = '\033[31m'
GREEN = '\033[32m'
YELLOW = '\033[33m'
BLUE = '\033[34m'
GREY = '\033[37m'
PURPLE = '\033[35m'
CYAN = '\033[36m'
ORANGE = '\033[38;5;208m'
PINK = '\033[38;5;201m'
LIGHT_BLUE = '\033[38;5;123m'
# Cool ANSI escape codes
# BLINK is an ANSI escape code for blinking text
# Example usage: print(f"{BLINK}This text will blink{RESET}")
BLINK = '\033[5m'
UNDERLINE = '\033[4m'
REVERSE = '\033[7m'
# PPT_BLUE = f'\033[38;5;{ansi_code}m'

# Define a color mapping
COLOR_MAP = {
    "white": (255, 255, 255),
    "black": (0, 0, 0),
    "red": (255, 0, 0),
    "green": (0, 255, 0),
    "blue": (0, 0, 255),
    "yellow": (255, 255, 0),
    "purple": (128, 0, 128),
    "default": (49, 51, 158),
}

def set_background_image(prs, slide, image_path, transparency=0.5):
    # Open the image and apply transparency
    img = Image.open(image_path).convert("RGBA")
    alpha = img.split()[3]
    alpha = alpha.point(lambda p: p * transparency)
    img.putalpha(alpha)
    
    # Save the modified image to a temporary file
    temp_image_path = "temp_image.png"
    img.save(temp_image_path)

    try:
        # Add the background image to the slide
        left = top = 0
        pic = slide.shapes.add_picture(temp_image_path, left, top, width=prs.slide_width, height=prs.slide_height)
        
        # Lock the picture
        pic_element = pic._element
        pic_element.set('noGrp', '1')  # Lock against grouping
        pic_element.set('noSelect', '1')  # Lock against selection
        pic_element.set('noRot', '1')  # Lock against rotation
        pic_element.set('noChangeAspect', '1')  # Lock against aspect ratio change
        pic_element.set('noMove', '1')  # Lock against moving
        pic_element.set('noResize', '1')  # Lock against resizing
        
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
    finally:
        # Remove the temporary image file
        os.remove(temp_image_path)

def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*COLOR_MAP[color])

def add_text_shadow(run):
    shadow_xml = """
    <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
        <a:outerShdw blurRad="38100" dist="38100" dir="5400000" algn="ctr" rotWithShape="0">
            <a:srgbClr val="000000">
                <a:alpha val="50000"/>
            </a:srgbClr>
        </a:outerShdw>
    </a:effectLst>
    """
    shadow_element = parse_xml(shadow_xml)
    run_element = run._r
    run_properties = run_element.get_or_add_rPr()
    run_properties.append(shadow_element)

def generate_ppt_from_txt(txt_file, ppt_file, background_image_path, fontsize, fontcolor, transparency, bgcolor):
    # Create a presentation object
    prs = Presentation()

    # Read the text file with UTF-8 encoding
    with open(txt_file, 'r', encoding='utf-8') as file:
        lines = file.readlines()

    if not lines:
        print("The text file is empty.")
        return

    title = lines[0].strip()  # The first line is the title for all slides
    content = ''.join(lines[1:])  # The rest of the file is the content

    # Replace \x0b with \n
    content = content.replace('\x0b', '\n')

    # Split by double newline and remove empty lines
    slides_content = [slide.strip() for slide in content.split('\n\n') if slide.strip()]

    left = top = 0
    # pic = slide.shapes.add_picture(background_image_path, left-0.1*prs.slide_width, top, height = prs.slide_height)


    for i, slide_content in enumerate(slides_content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use the title and content layout

        # Set the background image
        # set_background_image(slide, background_image_path)

        title_shape = slide.shapes.title
        title_shape.text = ""  # Clear any existing text

        # Add the title with underline
        p = title_shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(fontsize)
        run.font.bold = True
        run.font.underline = True
        run.font.color.rgb = RGBColor(*COLOR_MAP[fontcolor])
        add_text_shadow(run)  # Add shadow to each run in the paragraph

        # Add the slide number with smaller font size
        run = p.add_run()
        run.text = f" ({i+1}/{len(slides_content)})"
        run.font.size = Pt(fontsize * 0.3)
        run.font.color.rgb = RGBColor(*COLOR_MAP[fontcolor])
        add_text_shadow(run)  # Add shadow to each run in the paragraph

        # Get the text frame for the content
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()  # Clear any default content

        # Add content line by line
        for line in slide_content.split('\n'):
            stripped_line = line.strip()
            if stripped_line:  # Skip empty lines
                p = text_frame.add_paragraph()
                p.text = stripped_line  # Remove spaces at the end of the line

        if text_frame.text and text_frame.text[0] == '\n':
            text_frame.text = text_frame.text[1:]
            for i in range(len(text_frame.paragraphs)):
                text_frame.paragraphs[i]._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"))
                text_frame.paragraphs[i].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[i].font.size = Pt(fontsize)
                text_frame.paragraphs[i].font.bold = True
                text_frame.paragraphs[i].font.color.rgb = RGBColor(*COLOR_MAP[fontcolor])
                for run in text_frame.paragraphs[i].runs:
                    add_text_shadow(run)  # Add shadow to each run in the paragraph

        if background_image_path:
            set_background_image(prs, slide, background_image_path, transparency)
        else:
            set_background_color(slide, bgcolor)

    # Use title as ppt_file name if ppt_file is empty
    if not ppt_file:
        ppt_file = f"{title}.pptx"
    elif not ppt_file.lower().endswith('.pptx'):
        ppt_file += '.pptx'

    # Get the absolute path of the ppt_file
    ppt_file_path = os.path.abspath(ppt_file)

    try:
        # Add debug log before saving the presentation
        prs.save(ppt_file)
        print(f"{GREEN}{BOLD} >> Successfully saved PowerPoint presentation to {BLUE}{BOLD}{ppt_file_path}{RESET}")
    except Exception as e:
        print(f"Error: Failed to save PowerPoint presentation. {e}")
        print(f"{RED}{BOLD}Error: Failed to save PowerPoint presentation. {e}{RESET}")

# Enable virtual terminal processing on Windows
def enable_virtual_terminal_processing():
    kernel32 = ctypes.windll.kernel32
    handle = kernel32.GetStdHandle(-11)  # STD_OUTPUT_HANDLE
    mode = ctypes.c_ulong()
    kernel32.GetConsoleMode(handle, ctypes.byref(mode))
    mode.value |= 0x0004  # ENABLE_VIRTUAL_TERMINAL_PROCESSING
    kernel32.SetConsoleMode(handle, mode)

if __name__ == "__main__":

    # Call the function to enable virtual terminal processing
    if os.name == 'nt':
        enable_virtual_terminal_processing()

    # Welcome page with ASCII art
    ascii_art = """
   __         _         ___            __ 
  / /_ ______(_)______ |_  |___  ___  / /_
 / / // / __/ / __(_-</ __// _ \/ _ \/ __/
/_/\_, /_/ /_/\__/___/____/ .__/ .__/\__/ 
  /___/                  /_/  /_/         
"""
    print(f"{YELLOW}{BOLD}{ascii_art}{RESET}")
    print(f"{YELLOW}{REVERSE}Welcome to the PowerPoint Generator!{RESET}\n")
    
    print(f"{ORANGE}範例: 複製下面的命令並貼上到命令提示字元。{RESET}")
    print(f"{ORANGE}{BOLD}lyrics2ppt.exe --input input.txt{RESET}\n")

    # parser = argparse.ArgumentParser(description='Generate a PowerPoint presentation from a text file.')
    parser = argparse.ArgumentParser(
        description=ascii_art 
        + '\nGenerate a PowerPoint presentation from a text file.\n\n************\nMaintainer: Astra <astralee95@gmail.com>\n************\n\n'
        + "範例: \n lyrics2ppt.exe --input input.txt\n複製並貼上上面的命令。\n",
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument('--input',        type=str,   default="input.txt", help='Path to the input text file (輸入歌詞 txt 檔)')
    parser.add_argument('--output',       type=str,   default="", help='Path to the output PowerPoint file (輸出 ppt 檔名)')
    parser.add_argument('--bg-image',     type=str,   default="", help='Path to the background image file (default: none) (背景圖片文件的路徑, 預設: 無)')
    parser.add_argument('--bg-color',     type=str,   default='default', help='Background color (default: blue-purple) (背景顏色, 預設: 藍紫色)')
    parser.add_argument('--font-color',   type=str,   default='white', help='Font color (default: white) (標題 & 內文字型顏色, 預設: 白色)')
    parser.add_argument('--font-size',    type=float, default=48, help='Font size (default: 48) (標題 & 內文字型大小, 預設: 48)')
    parser.add_argument('--transparency', type=float, default=0.5, help='Transparency level for the background image (default: 0.5) (背景圖片的透明度, 預設: 0.5)')

    args = parser.parse_args()

    # Print the input arguments to the console
    print(f"{GREY}*********************")
    print(f"Input file      : {ITALIC}{args.input:<12}{RESET} (必須提供歌詞 txt 檔)")

    # Read the first line from the input file
    with open(args.input, 'r', encoding='utf-8') as file:
        first_line = file.readline().strip()
        first_line += '.pptx'

    if not args.output:
        print(f"Output file     : {ITALIC}{first_line}{RESET} (如果沒有提供, 則使用歌名為 ppt 檔名)")
    else:
        print(f"Output file     : {ITALIC}{args.output:<12}{RESET} (如果沒有提供, 則使用歌名為 ppt 檔名)")
    print("")
    
    # Convert background color to ANSI escape code
    ansi_bg_code = rgb_to_ansi(*COLOR_MAP[args.bg_color])
    PPT_BG_COLOR = f'\033[48;5;{ansi_bg_code}m'

    # Convert font color to ANSI escape code
    ansi_fg_code = rgb_to_ansi(*COLOR_MAP[args.font_color])
    FONT_COLOR = f'\033[38;5;{ansi_fg_code}m'
    
    print(f"Background color: {PPT_BG_COLOR}{ITALIC}{args.bg_color:<12}{RESET} (預設: ppt常用背景藍紫色, 選項: white, black, red, green, blue, yellow)")
    print(f"Font color      : {FONT_COLOR}{ITALIC}{args.font_color:<12}{RESET} (預設:白色, 選項: white, black, red, green, blue)")

    print(f"Font preview    : \033[48;5;{ansi_bg_code}m\033[38;5;{ansi_fg_code}m{ITALIC}{'PREVIEW':<12}{RESET} (預設:白色, 選項: white, black, red, green, blue)")
    print("")

    print(f"Font size       : {ITALIC}{args.font_size:<12}{RESET} (預設:48, 標題 & 內文字型大小)")
    print("")
    
    if args.bg_image:
        print(f"Background image: {ITALIC}{args.bg_image:<12}{RESET} (預設: 無)")
    else:
        print(f"Background image: {ITALIC}{'X':<12}{RESET} (預設: 無)")

    print(f"Transparency    : {ITALIC}{args.transparency:<12}{RESET} (如果有提供背景圖片, 則設定背景透明度)")
    print(f"*********************\n{RESET}")

    # if args.font_color not in COLOR_MAP:
    #     print(f"{RED}{BOLD}Error: Unsupported font color '{args.font_color}'. Supported colors are: {', '.join(COLOR_MAP.keys())}{RESET}")
    # elif args.bg_color not in COLOR_MAP:
    #     print(f"{RED}{BOLD}Error: Unsupported background color '{args.bg_color}'. Supported colors are: {', '.join(COLOR_MAP.keys())}{RESET}")
    # else:
    #     print(f"{YELLOW}{BOLD}Generating PowerPoint presentation...{RESET}\n")
    #     generate_ppt_from_txt(args.input, args.output, args.bg_image, args.font_size, args.font_color, args.transparency, args.bg_color)
    #     print(f"{YELLOW}{BOLD}\nPowerPoint presentation generation completed.{RESET}")

    if args.font_color not in COLOR_MAP:
        print(f"{RED}{BOLD}錯誤: 不支援的字型顏色 '{args.font_color}'。支援的顏色有: {', '.join(COLOR_MAP.keys())}{RESET}")
    elif args.bg_color not in COLOR_MAP:
        print(f"{RED}{BOLD}錯誤: 不支援的背景顏色 '{args.bg_color}'。支援的顏色有: {', '.join(COLOR_MAP.keys())}{RESET}")
    else:
        print(f"{YELLOW}{BOLD}正在生成 PowerPoint 簡報...{RESET}\n")
        generate_ppt_from_txt(args.input, args.output, args.bg_image, args.font_size, args.font_color, args.transparency, args.bg_color)
        print(f"{YELLOW}{BOLD}\nPowerPoint 簡報生成完成。{RESET}")
